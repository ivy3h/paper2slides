#!/usr/bin/env python
"""Build the shipped example and assert the things that have actually broken before.

Usage: python selftest.py [deck.json]

No test framework, no fixtures: it builds a real deck into a temp file and reads
the PPTX back. Every assertion here corresponds to a bug that shipped at least once.
"""
import sys, os, json, tempfile, re
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
import build_slides as B

DASH_RE = re.compile("[‐‑‒–—―−⁃﹘﹣－]")

failures = []


def check(cond, msg):
    print(("  ok    " if cond else "  FAIL  ") + msg)
    if not cond:
        failures.append(msg)


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "..", "examples", "spurious-rewards.deck.json")
    deck = json.load(open(deck_path, encoding="utf-8"))
    out = tempfile.mktemp(suffix=".pptx")
    B.build(deck, out, os.path.join(os.path.dirname(deck_path)))
    prs = Presentation(out)
    slides = list(prs.slides)

    print("\nstructure")
    check(len(slides) == len(deck["slides"]), f"slide count matches deck.json ({len(slides)})")
    check(abs(prs.slide_width / 914400 - 13.333) < 0.01 and abs(prs.slide_height / 914400 - 7.5) < 0.01,
          "16:9 at 13.333 x 7.5 in")

    body, sizes, supers = [], set(), 0
    for s in slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            body.append(sh.text_frame.text)
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        sizes.add(round(r.font.size / 12700, 2))
                    if r.font._rPr.get("baseline"):
                        supers += 1
        if s.has_notes_slide:
            body.append(s.notes_slide.notes_text_frame.text)
    blob = "\n".join(body)

    print("\ntext")
    check(not DASH_RE.search(blob), "no em/en dash survives into the PPTX")
    check("[^" not in blob, "no citation marker leaks through unrendered")
    check("**" not in blob, "no bold markup leaks through unrendered")
    check(supers > 0, f"superscript markers render as superscripts ({supers} runs)")

    print("\ntypography")
    # the house rule is that body copy is one size everywhere; titles and chrome sit
    # outside that band, and the footer credit is allowed to shrink to stay on one line
    body_band = sorted(x for x in sizes if 16 <= x <= 26)
    check(body_band == [B.BODY_SIZE],
          f"body copy is exactly one size ({B.BODY_SIZE}pt), found {body_band}")
    check(min(sizes) >= 6, f"nothing renders below 6pt (smallest {min(sizes)}pt)")

    print("\nfooter")
    # the rule sits at RULE_Y; footer text must start below it or it gets struck through
    check(B.FOOT_Y > B.RULE_Y + 0.022, "footer text clears the gold rule")
    check(B.FOOT_Y + 0.17 <= 7.5, "footer text stays on the slide")
    numbered = 0
    for i, s in enumerate(slides, 1):
        want = f"{i} / {len(slides)}"
        if any(sh.has_text_frame and sh.text_frame.text.strip() == want for sh in s.shapes):
            numbered += 1
    check(numbered == len(slides) - 1, f"every content slide numbered correctly ({numbered})")

    print("\nnotes")
    thin = [i for i, s in enumerate(slides, 1)
            if not s.has_notes_slide or len(s.notes_slide.notes_text_frame.text.strip()) < 40]
    check(not thin, f"every slide has speaker notes {thin if thin else ''}")

    print("\nlayout fit")
    over = B.OVERFLOWS
    check(not over, f"no text block overflows its box {over if over else ''}")

    os.unlink(out)
    print(f"\n{len(failures)} failure(s)")
    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
