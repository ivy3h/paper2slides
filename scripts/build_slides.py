#!/usr/bin/env python
"""Build a 16:9 editable PPTX from a deck.json spec.

Usage: python build_slides.py deck.json slides.pptx [assets_dir]

deck.json schema (see examples/spurious-rewards.deck.json):
{
  "meta": {title, authors, affiliation, venue, presenter, date},
  "slides": [ { "layout": <type>, "title":..., "kicker":..., "number":...,
                "bullets": [str | {"text":str,"level":0}], "columns": [[..],[..]],
                "figure": "<manifest id or filename>", "figure_caption": str,
                "cite": str | [str],   # source credit, footer-left beside the page number
                                       # a list auto-numbers; mark the body with [^1], [^2], [^1,2]
                "notes": str } ]
}
Layouts: title | section | bullets | bullets_figure | figure | figure_bullets | table |
         two_column | matrix | takeaways
Bold spans inside text use **markdown** style: "**Random rewards** still gain +21.4".
House style: no em/en dashes in deck text; use commas, colons or parentheses instead.
Body text is one size (BODY_SIZE) on every content layout; no title-slide subtitle.
"""
import sys, json, os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

NAVY  = RGBColor(0x0B, 0x25, 0x45)
GOLD  = RGBColor(0xE8, 0xA3, 0x3D)
GOLDD = RGBColor(0xB9, 0x78, 0x1B)
INK   = RGBColor(0x1A, 0x24, 0x38)
MUTED = RGBColor(0x5A, 0x66, 0x80)
GOLDSOFT = RGBColor(0xFB, 0xF1, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Helvetica Neue"

SW, SH = Inches(13.333), Inches(7.5)

def parse_bold(s):
    """Split '**bold** normal' into [(text, is_bold), ...]."""
    parts, bold = [], False
    for chunk in s.split("**"):
        if chunk:
            parts.append((chunk, bold))
        bold = not bold
    return parts or [("", False)]

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf

SUP_RE = re.compile(r"\[\^([\d,]+)\]")  # [^1] or [^1,2] -> superscript, tying body text to a footer cite

def _run(para, text, size, color, bold, italic, font, sup=False):
    r = para.add_run(); r.text = text
    r.font.size = Pt(size * 0.72 if sup else size); r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    if sup:
        r.font._rPr.set("baseline", "30000")  # raise 30%, DrawingML's superscript
    return r

def style_runs(para, text, size, color, bold=False, italic=False, font=FONT):
    """Add runs to a paragraph, honoring **bold** spans (bold+navy) and [^N] superscripts."""
    for chunk, is_b in parse_bold(text):
        col = NAVY if (is_b and not bold) else color
        pos = 0
        for m in SUP_RE.finditer(chunk):
            if m.start() > pos:
                _run(para, chunk[pos:m.start()], size, col, bold or is_b, italic, font)
            _run(para, m.group(1), size, col, bold or is_b, italic, font, sup=True)
            pos = m.end()
        if pos < len(chunk):
            _run(para, chunk[pos:], size, col, bold or is_b, italic, font)

def add_para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
             space_after=6, space_before=0, line=1.04, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    try: p.line_spacing = line
    except Exception: pass
    style_runs(p, text, size, color, bold, italic)
    return p

RULE_Y = 7.18          # gold hairline above the footer
FOOT_Y = RULE_Y + 0.08  # text sits clear of the rule (a 10pt line is ~0.17in tall)

def footer(slide, num, total, cite=None):
    """Gold hairline, optional source credit bottom-left, page number bottom-right."""
    rect(slide, 0, RULE_Y, 13.333, 0.022, GOLD)
    if cite:
        if isinstance(cite, (list, tuple)):
            # auto-number, so [^1] in the body lands on the first entry here
            cite = "   ·   ".join(f"[^{i}] {c}" for i, c in enumerate(cite, 1))
        cite = str(cite)
        # shrink rather than wrap off the bottom of the slide
        n = len(SUP_RE.sub("x", cite))
        fits = chars_per_line(10.2, 10)
        size = 10 if n <= fits else max(7.5, round(10 * fits / n, 1))
        tb, tf = textbox(slide, 0.55, FOOT_Y, 10.2, 0.24)  # stops short of the page number
        add_para(tf, cite, size, MUTED, first=True)
    tb, tf = textbox(slide, 11.0, FOOT_Y, 1.78, 0.24)
    add_para(tf, f"{num} / {total}", 10, MUTED, align=PP_ALIGN.RIGHT, first=True)

def add_notes(slide, notes):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def fit(path, box_w, box_h):
    """Return (w,h) in inches fitting image inside box preserving aspect."""
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return box_w, box_h
    ar = iw / ih
    w, h = box_w, box_w / ar
    if h > box_h:
        h, w = box_h, box_h * ar
    return w, h

def add_figure(slide, path, cx, cy, box_w, box_h, caption=None, cap_h=0.45, cap_size=11.5):
    """Place image centered in box (cx,cy = box top-left), caption below."""
    if not path or not os.path.exists(path):
        return
    reserve = cap_h if caption else 0
    w, h = fit(path, box_w, box_h - reserve)
    left = cx + (box_w - w) / 2
    top = cy + (box_h - reserve - h) / 2
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    pic.line.color.rgb = RGBColor(0xD7, 0xDE, 0xEA); pic.line.width = Pt(0.75)
    if caption:
        # keep the caption attached to the image instead of pinned to the bottom of the box
        cap_top = min(top + h + 0.06, cy + box_h - cap_h)
        tb, tf = textbox(slide, cx, cap_top, box_w, cap_h)
        add_para(tf, caption, cap_size, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True, line=1.05)

def slide_title(slide, title, kicker=None):
    if kicker:
        tb, tf = textbox(slide, 0.55, 0.34, 12.2, 0.32)
        p = add_para(tf, kicker.upper(), 12.5, GOLDD, bold=True, first=True)
    tb, tf = textbox(slide, 0.55, 0.62 if kicker else 0.45, 12.2, 0.9)
    add_para(tf, title, 27, NAVY, bold=True, first=True, line=1.0)
    rect(slide, 0.57, 1.45 if kicker else 1.28, 1.6, 0.045, GOLD)
    return 1.75 if kicker else 1.58  # content top y

def hanging(para, indent_in):
    """Hanging indent: wrapped lines align under the text, not under the bullet glyph."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(indent_in))))
    pPr.set("indent", str(-int(Inches(indent_in))))

BODY_SIZE = 17  # one body size for every content layout, so slides don't jump around

# Measured on rendered output: Helvetica Neue mixed-case advances ~0.00616 in per
# point of font size. Everything that needs to guess how much room text takes
# (bullet overflow, the takeaways box, the footer credit) derives from this one
# number rather than carrying its own magic constant.
CHAR_W_PT = 0.00616

def chars_per_line(width_in, size_pt):
    return max(1, int(width_in / (CHAR_W_PT * size_pt)))

def text_lines(text, width_in, size_pt):
    text = str(text).replace("**", "")
    text = SUP_RE.sub("x", text)
    return max(1, -(-len(text) // chars_per_line(width_in, size_pt)))

def block_height(bullets, width_in, size_pt, pad_pt=9):
    """Rough height in inches of a bullets_block, matching its spacing."""
    total = 0.0
    for b in bullets:
        text, level = (b.get("text", ""), b.get("level", 0)) if isinstance(b, dict) else (b, 0)
        sz = size_pt if level == 0 else size_pt - 2
        indent = 0.6 if level else 0.0
        total += text_lines(text, width_in - indent, sz) * sz * 1.06 / 72 + pad_pt / 72
    return total

OVERFLOWS = []  # (where, needed_in, available_in), reported at the end of build()

def bullets_block(slide, bullets, x, y, w, h, base_size=BODY_SIZE, where=None):
    n = len(bullets)
    size = base_size if n <= 9 else 15  # only shrink if a slide is unusually dense
    need = block_height(bullets, w, size)
    if need > h + 0.02:
        OVERFLOWS.append((where or "?", round(need, 2), round(h, 2)))
    tb, tf = textbox(slide, x, y, w, h)
    first = True
    for b in bullets:
        if isinstance(b, dict):
            text, level = b.get("text", ""), b.get("level", 0)
        else:
            text, level = str(b), 0
        if level == 0:
            p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
            p.space_after = Pt(7); p.space_before = Pt(2); p.line_spacing = 1.06
            r = p.add_run(); r.text = "▪  "; r.font.size = Pt(size); r.font.color.rgb = GOLD; r.font.bold = True; r.font.name = FONT
            style_runs(p, text, size, INK)
            hanging(p, size / 68.0)  # ≈ width of "▪  " at this point size
        else:
            p = tf.add_paragraph(); p.space_after = Pt(4); p.line_spacing = 1.04
            r = p.add_run(); r.text = "        ·  "; r.font.size = Pt(size-2); r.font.color.rgb = MUTED; r.font.name = FONT
            style_runs(p, text, size-2, INK)
            hanging(p, (size - 2) / 26.0)  # ≈ width of the indented sub-bullet prefix
        first = False
    return tb

def build(deck, out_pptx, assets):
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]
    meta = deck.get("meta", {})
    slides = deck["slides"]
    total = len(slides)

    def figpath(fid):
        if not fid: return None
        for cand in (fid, fid + ".png", os.path.join(assets, fid), os.path.join(assets, fid + ".png")):
            if os.path.exists(cand): return cand
        return None

    for i, s in enumerate(slides, 1):
        layout = s.get("layout", "bullets")
        sl = prs.slides.add_slide(blank)
        notes = s.get("notes", "")

        if layout == "title":
            set_bg(sl, NAVY)
            rect(sl, 0, 0, 13.333, 0.16, GOLD)
            tb, tf = textbox(sl, 0.9, 0.5, 11.5, 0.4)
            add_para(tf, (meta.get("venue") or "").upper(), 14, GOLD, bold=True, first=True)
            # centred between the venue strip and the gold rule, so 1, 2 or 3 line
            # titles all sit optically right without a subtitle to pad them out
            tb, tf = textbox(sl, 0.9, 1.15, 11.5, 3.3, anchor=MSO_ANCHOR.MIDDLE)
            add_para(tf, meta.get("title", ""), 40, WHITE, bold=True, first=True, line=1.03)
            rect(sl, 0.92, 4.7, 2.2, 0.05, GOLD)
            tb, tf = textbox(sl, 0.9, 4.95, 11.5, 1.6)
            add_para(tf, meta.get("authors", ""), 17, RGBColor(0xDF,0xE7,0xF4), bold=True, first=True)
            if meta.get("affiliation"):
                add_para(tf, meta["affiliation"], 14, RGBColor(0xAA,0xB9,0xD4), italic=True, space_before=2)
            foot = "   ·   ".join(x for x in [meta.get("presenter",""), meta.get("date","")] if x)
            if foot:
                add_para(tf, foot, 13.5, GOLD, space_before=14)
            add_notes(sl, notes); continue

        if layout == "section":
            set_bg(sl, NAVY)
            rect(sl, 0, 0, 0.22, 7.5, GOLD)
            tb, tf = textbox(sl, 1.1, 2.4, 11, 2.6, anchor=MSO_ANCHOR.MIDDLE)
            if s.get("number"):
                add_para(tf, str(s["number"]), 60, GOLD, bold=True, first=True, line=1.0)
            add_para(tf, s.get("title",""), 36, WHITE, bold=True, space_before=4, line=1.02)
            if s.get("subtitle"):
                add_para(tf, s["subtitle"], 18, RGBColor(0xC8,0xD4,0xEA), space_before=8)
            add_notes(sl, notes); continue

        # ---- content slides (white) ----
        set_bg(sl, WHITE)
        cy = slide_title(sl, s.get("title",""), s.get("kicker"))

        if layout == "bullets":
            bullets_block(sl, s.get("bullets", []), 0.6, cy+0.05, 12.1, 5.1, where=i)

        elif layout == "bullets_figure":
            bullets_block(sl, s.get("bullets", []), 0.6, cy+0.05, 6.6, 5.0, where=i)
            add_figure(sl, figpath(s.get("figure")), 7.45, cy+0.0, 5.3, 5.05, s.get("figure_caption"))

        elif layout == "figure":
            add_figure(sl, figpath(s.get("figure")), 0.7, cy+0.0, 11.9, 5.05, s.get("figure_caption"))

        elif layout == "figure_bullets":
            # wide, short figures (tables / multi-panel strips) waste a full-slide `figure`:
            # give them a full-width band on top and put the takeaways underneath.
            fh = float(s.get("figure_height", 3.2))
            add_figure(sl, figpath(s.get("figure")), 0.6, cy+0.0, 12.13, fh,
                       s.get("figure_caption"), cap_h=0.32, cap_size=10.5)
            rest = 5.05 - fh - 0.15
            bullets_block(sl, s.get("bullets", []), 0.6, cy + fh + 0.15, 12.13, max(rest, 0.6),
                          base_size=float(s.get("bullet_size", BODY_SIZE)), where=i)

        elif layout == "table":
            if s.get("bullets"):
                bullets_block(sl, s["bullets"], 0.6, cy+0.05, 4.6, 5.0, where=i)
                add_figure(sl, figpath(s.get("figure")), 5.4, cy+0.0, 7.4, 5.05, s.get("figure_caption"))
            else:
                # full-bleed results table: maximize footprint so digits stay legible.
                box_top = 1.55  # start just under the title's gold rule
                box_h = 7.13 - box_top  # extend down to just above the footer rule
                add_figure(sl, figpath(s.get("figure")), 0.3, box_top, 12.73, box_h,
                           s.get("figure_caption"), cap_h=0.3, cap_size=10.5)

        elif layout == "matrix":
            # case-analysis grid (native shapes, still editable): row/col headers + cells,
            # cells listed in "dead" render muted (the branch that gets killed).
            m = s.get("matrix", {})
            mcols, mrows = m.get("cols", []), m.get("rows", [])
            mcells = m.get("cells", [])
            deadset = {tuple(d) for d in m.get("dead", [])}
            gx, gy, gw = 0.6, cy + 0.05, 12.13
            lab_w = float(m.get("label_width", 3.4))
            grid_h = float(m.get("height", 2.9))
            hh = 0.5
            cw = (gw - lab_w) / max(len(mcols), 1)
            rh = (grid_h - hh) / max(len(mrows), 1)
            for cj, ch in enumerate(mcols):  # NB: do not shadow the outer slide counter `i`
                tb, tf = textbox(sl, gx + lab_w + cj * cw + 0.1, gy, cw - 0.2, hh, anchor=MSO_ANCHOR.MIDDLE)
                add_para(tf, ch, 14, GOLDD, bold=True, align=PP_ALIGN.CENTER, first=True, line=1.05)
            for ri, rhd in enumerate(mrows):
                top = gy + hh + ri * rh
                tb, tf = textbox(sl, gx, top + 0.08, lab_w - 0.25, rh - 0.16, anchor=MSO_ANCHOR.MIDDLE)
                add_para(tf, rhd, 15, NAVY, bold=True, first=True, line=1.06)
                for cj in range(len(mcols)):
                    is_dead = (ri, cj) in deadset
                    rect(sl, gx + lab_w + cj * cw + 0.06, top + 0.06, cw - 0.12, rh - 0.12,
                         RGBColor(0xEF, 0xF1, 0xF5) if is_dead else GOLDSOFT,
                         line=RGBColor(0xD7, 0xDE, 0xEA) if is_dead else RGBColor(0xEC, 0xCF, 0x97))
                    txt = mcells[ri][cj] if ri < len(mcells) and cj < len(mcells[ri]) else ""
                    tb, tf = textbox(sl, gx + lab_w + cj * cw + 0.24, top + 0.14, cw - 0.48, rh - 0.28,
                                     anchor=MSO_ANCHOR.MIDDLE)
                    add_para(tf, txt, 15, MUTED if is_dead else INK, align=PP_ALIGN.CENTER,
                             first=True, line=1.08, italic=is_dead)
            if s.get("bullets"):
                bullets_block(sl, s["bullets"], 0.6, gy + grid_h + 0.22, 12.13,
                              max(5.05 - grid_h - 0.27, 0.6), where=i)

        elif layout == "two_column":
            cols = s.get("columns", [[], []])
            bullets_block(sl, cols[0] if len(cols)>0 else [], 0.6, cy+0.05, 5.9, 5.0, where=i)
            rect(sl, 6.66, cy+0.1, 0.012, 4.7, RGBColor(0xD7,0xDE,0xEA))
            bullets_block(sl, cols[1] if len(cols)>1 else [], 6.95, cy+0.05, 5.8, 5.0, where=i)

        elif layout == "takeaways":
            # size the callout to its content so the box hugs the text at any bullet count
            bl = s.get("bullets", [])
            box_h = min(4.6, max(1.5, 0.25 + block_height(bl, 11.4, BODY_SIZE) + 0.25))
            box = rect(sl, 0.6, cy+0.2, 12.13, box_h, GOLDSOFT)
            box.line.color.rgb = RGBColor(0xEC,0xCF,0x97); box.line.width = Pt(1)
            rect(sl, 0.6, cy+0.2, 0.13, box_h, GOLD)
            bullets_block(sl, bl, 1.05, cy+0.45, 11.4, box_h - 0.5, where=i)

        else:
            bullets_block(sl, s.get("bullets", []), 0.6, cy+0.05, 12.1, 5.1, where=i)

        footer(sl, i, total, s.get("cite"))
        add_notes(sl, notes)

    prs.save(out_pptx)
    print(f"saved {out_pptx}  ({total} slides)")
    for where, need, have in OVERFLOWS:
        print(f"  warning: slide {where} text needs ~{need}in in a {have}in box", file=sys.stderr)

if __name__ == "__main__":
    deck_path = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else "slides.pptx"
    assets = sys.argv[3] if len(sys.argv) > 3 else "assets"
    deck = json.load(open(deck_path))
    build(deck, out, assets)
